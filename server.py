import asyncio
import os
from pathlib import Path
from zipfile import ZipFile
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from mcp.server.fastmcp import FastMCP

# Initialize MCP server
mcp = FastMCP("office-mcp")

# Base directory for documents
DOCS_DIR = Path("./documents")
DOCS_DIR.mkdir(exist_ok=True)

def get_path(filename: str, ext: str) -> Path:
    if not filename.endswith(ext):
        filename += ext
    return DOCS_DIR / filename

# -------------------------------
# WORD TOOLS (.docx)
# -------------------------------

@mcp.tool()
def create_docx(filename: str, content: str = "") -> str:
    path = get_path(filename, ".docx")
    if path.exists():
        return f"❌ File already exists: {path}"

    try:
        doc = Document()
        doc.add_paragraph(content if content.strip() else "(new document)")
        doc.save(path)
        return f"✅ Created Word document: {path}"
    except Exception as e:
        return f"❌ Failed to create Word doc: {str(e)}"

@mcp.tool()
def read_docx(filename: str) -> str:
    path = get_path(filename, ".docx")
    if not path.exists():
        return f"❌ File not found: {path}"

    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs) or "(empty document)"
    except Exception as e:
        try:
            with ZipFile(path, "r") as zf:
                files = zf.namelist()
                if "word/document.xml" not in files:
                    return f"❌ Not a Word file. Contains: {files[:10]}..."
        except Exception:
            pass
        return f"❌ Failed to read Word doc: {str(e)}"

@mcp.tool()
def update_docx(filename: str, changes: list[str]) -> str:
    path = get_path(filename, ".docx")
    if not path.exists():
        return f"❌ File not found: {path}"

    try:
        doc = Document(path)
        for change in changes:
            doc.add_paragraph(change)
        doc.save(path)
        return f"✅ Updated Word doc {filename} with {len(changes)} changes."
    except Exception as e:
        return f"❌ Failed to update Word doc: {str(e)}"

@mcp.tool()
def delete_docx(filename: str) -> str:
    path = get_path(filename, ".docx")
    if not path.exists():
        return f"❌ File not found: {path}"
    try:
        os.remove(path)
        return f"🗑️ Deleted Word doc: {path}"
    except Exception as e:
        return f"❌ Failed to delete Word doc: {str(e)}"

# -------------------------------
# EXCEL TOOLS (.xlsx)
# -------------------------------

@mcp.tool()
def create_xlsx(filename: str, sheet_name: str = "Sheet1", data: list[list[str]] = None) -> str:
    path = get_path(filename, ".xlsx")
    if path.exists():
        return f"❌ File already exists: {path}"

    try:
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        if data:
            for row in data:
                ws.append(row)
        wb.save(path)
        return f"✅ Created Excel file: {path}"
    except Exception as e:
        return f"❌ Failed to create Excel: {str(e)}"

@mcp.tool()
def read_xlsx(filename: str, sheet_name: str = None) -> str:
    path = get_path(filename, ".xlsx")
    if not path.exists():
        return f"❌ File not found: {path}"

    try:
        wb = load_workbook(path, data_only=True)
        ws = wb[sheet_name] if sheet_name else wb.active
        rows = [[str(cell.value) if cell.value is not None else "" for cell in row] for row in ws.iter_rows()]
        return "\n".join([", ".join(r) for r in rows]) or "(empty sheet)"
    except Exception as e:
        return f"❌ Failed to read Excel: {str(e)}"

@mcp.tool()
def update_xlsx(filename: str, sheet_name: str, data: list[list[str]]) -> str:
    path = get_path(filename, ".xlsx")
    if not path.exists():
        return f"❌ File not found: {path}"

    try:
        wb = load_workbook(path)
        ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.create_sheet(sheet_name)
        for row in data:
            ws.append(row)
        wb.save(path)
        return f"✅ Updated Excel file {filename} with {len(data)} rows."
    except Exception as e:
        return f"❌ Failed to update Excel: {str(e)}"

@mcp.tool()
def delete_xlsx(filename: str) -> str:
    path = get_path(filename, ".xlsx")
    if not path.exists():
        return f"❌ File not found: {path}"
    try:
        os.remove(path)
        return f"🗑️ Deleted Excel file: {path}"
    except Exception as e:
        return f"❌ Failed to delete Excel: {str(e)}"

# -------------------------------
# POWERPOINT TOOLS (.pptx)
# -------------------------------

@mcp.tool()
def create_pptx(filename: str, title: str = "New Presentation", content: str = "") -> str:
    path = get_path(filename, ".pptx")
    if path.exists():
        return f"❌ File already exists: {path}"

    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        prs.save(path)
        return f"✅ Created PowerPoint: {path}"
    except Exception as e:
        return f"❌ Failed to create PowerPoint: {str(e)}"

@mcp.tool()
def read_pptx(filename: str) -> str:
    path = get_path(filename, ".pptx")
    if not path.exists():
        return f"❌ File not found: {path}"

    try:
        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            slides_text.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(slides_text) or "(empty presentation)"
    except Exception as e:
        return f"❌ Failed to read PowerPoint: {str(e)}"

@mcp.tool()
def update_pptx(filename: str, slides: list[dict]) -> str:
    """
    slides = [{"title": "Slide Title", "content": "Some text"}]
    """
    path = get_path(filename, ".pptx")
    if not path.exists():
        return f"❌ File not found: {path}"

    try:
        prs = Presentation(path)
        for s in slides:
            layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s.get("title", "")
            slide.placeholders[1].text = s.get("content", "")
        prs.save(path)
        return f"✅ Updated PowerPoint {filename} with {len(slides)} new slides."
    except Exception as e:
        return f"❌ Failed to update PowerPoint: {str(e)}"

@mcp.tool()
def delete_pptx(filename: str) -> str:
    path = get_path(filename, ".pptx")
    if not path.exists():
        return f"❌ File not found: {path}"
    try:
        os.remove(path)
        return f"🗑️ Deleted PowerPoint: {path}"
    except Exception as e:
        return f"❌ Failed to delete PowerPoint: {str(e)}"

# -------------------------------
# MAIN ENTRY
# -------------------------------
if __name__ == "__main__":
    asyncio.run(mcp.run())
